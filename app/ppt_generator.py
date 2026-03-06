from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
import os


def generate_ppt(plan):
    os.makedirs("artifacts", exist_ok=True)
    prs = Presentation()

    for slide in plan["slides"]:

        if slide["type"] == "title":
            layout = prs.slide_layouts[0]
            s = prs.slides.add_slide(layout)
            s.shapes.title.text = slide["title"]
            s.placeholders[1].text = slide.get("subtitle","")

        elif slide["type"] == "bullet":
            layout = prs.slide_layouts[1]
            s = prs.slides.add_slide(layout)
            s.shapes.title.text = slide["title"]

            body = s.shapes.placeholders[1].text = slide["points"][0]

        elif slide["type"] == "chart":

            layout = prs.slide_layouts[5]
            s = prs.slides.add_slide(layout)

            s.shapes.title.text = slide["title"]

            chart_data = CategoryChartData()

            categories = []
            values = []

            for k, v in slide["data"].items():
                categories.append(str(k))
                values.append(float(v))

            chart_data.categories = categories

            # dynamic series name
            series_name = slide.get("title", "Series")

            chart_data.add_series(series_name, values)

            chart_type_map = {
                "line": XL_CHART_TYPE.LINE,
                "bar": XL_CHART_TYPE.BAR_CLUSTERED,
                "column": XL_CHART_TYPE.COLUMN_CLUSTERED
            }

            chart_type = chart_type_map.get(
                slide.get("chart_type"),
                XL_CHART_TYPE.COLUMN_CLUSTERED
            )

            s.shapes.add_chart(
                chart_type,
                0, 0, 5000000, 3000000,
                chart_data
            )

    prs.save("artifacts/generated_deck.pptx")